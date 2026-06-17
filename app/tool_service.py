"""
工具服务 - 真正可用的工具集
包含：天气、计算器、时间、记忆、文档检索、文件操作、打开应用、执行命令、网页搜索、图片分析、翻译、UI自动化、截图
"""
import ast
import base64
import io
import json
import logging
import operator
import os
import re
import sqlite3
import subprocess
import shutil
import tempfile
import time
from datetime import datetime, timezone, timedelta
from pathlib import Path
from typing import Any, Optional

import httpx

from app.config import UPLOAD_DIR, MEMORY_DB_PATH, LLM_API_KEY, LLM_BASE_URL, LLM_MODEL
from app.config import SILICONFLOW_API_KEY, SILICONFLOW_BASE_URL, SILICONFLOW_IMAGE_MODEL
from app.config import ZHIPU_API_KEY, ZHIPU_BASE_URL

# 通义千问视觉模型配置（用于图片分析）
DASHSCOPE_API_KEY = os.getenv("DASHSCOPE_API_KEY", "")
DASHSCOPE_VISION_URL = "https://dashscope.aliyuncs.com/compatible_mode/v1"
DASHSCOPE_VISION_MODEL = "qwen-vl-plus"

logger = logging.getLogger(__name__)

# 安全的数学运算符
SAFE_OPERATORS = {
    ast.Add: operator.add, ast.Sub: operator.sub,
    ast.Mult: operator.mul, ast.Div: operator.truediv,
    ast.FloorDiv: operator.floordiv, ast.Mod: operator.mod,
    ast.Pow: operator.pow, ast.USub: operator.neg, ast.UAdd: operator.pos,
}

# ==================== 城市坐标映射 ====================
CITY_COORDS = {
    "北京": (39.90, 116.40), "上海": (31.23, 121.47), "广州": (23.13, 113.26),
    "深圳": (22.54, 114.06), "成都": (30.57, 104.07), "杭州": (30.27, 120.15),
    "武汉": (30.59, 114.31), "南京": (32.06, 118.80), "西安": (34.26, 108.94),
    "重庆": (29.56, 106.55), "天津": (39.13, 117.20), "苏州": (31.30, 120.62),
    "长沙": (28.23, 112.94), "郑州": (34.75, 113.65), "青岛": (36.07, 120.38),
    "大连": (38.91, 121.60), "厦门": (24.48, 118.09), "昆明": (25.04, 102.68),
    "哈尔滨": (45.75, 126.65), "沈阳": (41.80, 123.43), "济南": (36.65, 116.98),
    "合肥": (31.82, 117.23), "福州": (26.07, 119.30), "南宁": (22.82, 108.37),
    "贵阳": (26.65, 106.63), "太原": (37.87, 112.55), "石家庄": (38.04, 114.51),
    "兰州": (36.06, 103.83), "乌鲁木齐": (43.83, 87.62), "拉萨": (29.65, 91.17),
    "海口": (20.04, 110.35), "三亚": (18.25, 109.50), "香港": (22.32, 114.17),
    "台北": (25.03, 121.57),
}

WEATHER_CODES = {
    0: "晴", 1: "大部晴", 2: "多云", 3: "阴天",
    45: "雾", 48: "雾凇", 51: "小毛毛雨", 53: "毛毛雨", 55: "大毛毛雨",
    61: "小雨", 63: "中雨", 65: "大雨", 66: "冻雨", 67: "大冻雨",
    71: "小雪", 73: "中雪", 75: "大雪", 77: "雪粒",
    80: "阵雨", 81: "中阵雨", 82: "大阵雨",
    85: "阵雪", 86: "大阵雪",
    95: "雷暴", 96: "雷暴+冰雹", 99: "强雷暴+冰雹",
}

# ==================== 系统保护路径（禁止删/改）====================
PROTECTED_PATHS = [
    r"c:\windows", r"c:\program files", r"c:\program files (x86)",
    r"c:\programdata", r"c:\users\all users",
]


def _safe_eval(expr: str) -> float:
    """安全数学表达式计算"""
    expr = expr.replace("×", "*").replace("÷", "/").replace("％", "%").replace("^", "**").strip()
    try:
        tree = ast.parse(expr, mode="eval")
    except SyntaxError:
        raise ValueError(f"无效表达式: {expr}")

    def _eval_node(node):
        if isinstance(node, ast.Expression):
            return _eval_node(node.body)
        elif isinstance(node, ast.Constant) and isinstance(node.value, (int, float)):
            return node.value
        elif isinstance(node, ast.Num):
            return node.n
        elif isinstance(node, ast.BinOp):
            op_type = type(node.op)
            if op_type not in SAFE_OPERATORS:
                raise ValueError(f"不支持的运算符: {op_type.__name__}")
            return SAFE_OPERATORS[op_type](_eval_node(node.left), _eval_node(node.right))
        elif isinstance(node, ast.UnaryOp):
            op_type = type(node.op)
            if op_type not in SAFE_OPERATORS:
                raise ValueError(f"不支持的运算符: {op_type.__name__}")
            return SAFE_OPERATORS[op_type](_eval_node(node.operand))
        else:
            raise ValueError(f"不支持的表达式: {type(node).__name__}")

    result = _eval_node(tree)
    return int(result) if isinstance(result, float) and result == int(result) else round(result, 6)


class ToolService:
    """工具执行服务 - 支持人机确认机制"""

    # 危险工具列表（需要用户确认才能执行）
    DANGEROUS_TOOLS = {"delete_file", "run_command", "app_action", "app_keyboard", "app_click_at", "app_drag"}

    def __init__(self, memory_service=None, rag_service=None):
        self.memory_service = memory_service
        self.rag_service = rag_service
        # LLM 客户端（用于图片分析和翻译）
        self._llm_client = None
        # 智谱 GLM 客户端（用于翻译，未配置时回退到主客户端）
        self._zhipu_client = None

    def set_llm_client(self, client):
        """设置 LangChain ChatOpenAI 客户端"""
        self._llm_client = client

    def set_zhipu_client(self, zhipu_client):
        """设置智谱 GLM 客户端（用于翻译）；为 None 时使用主客户端"""
        self._zhipu_client = zhipu_client or self._llm_client

    async def execute(
        self, tool_name: str, args: dict, user_id: str = "", confirmed: bool = False
    ) -> dict:
        """
        执行工具并返回结果。
        返回 {"content": str, "need_confirm": bool, "pending_id": str|None}
        """
        try:
            handler = getattr(self, f"_tool_{tool_name}", None)
            if handler is None:
                # 检查是否为 MCP 外部工具
                if tool_name.startswith("mcp_"):
                    return await self._execute_mcp_tool(tool_name, args)
                available = [
                    "get_weather", "calculator", "get_current_time",
                    "save_memory", "search_memory", "search_documents",
                    "list_directory", "read_file", "write_file",
                    "delete_file", "open_file", "run_command", "web_search",
                    "analyze_image", "translate", "generate_image",
                    "app_get_controls", "app_action", "app_keyboard",
                    "app_screenshot", "app_click_at", "app_clipboard",
                    "app_list_windows", "app_drag", "app_write_text",
                ]
                return {"content": f"未知工具: {tool_name}。可用: {', '.join(available)}", "need_confirm": False}

            # 危险操作 → 需要确认
            if tool_name in self.DANGEROUS_TOOLS and not confirmed:
                pending_id = self._save_pending_action(user_id, tool_name, args)
                preview = self._preview_dangerous_action(tool_name, args)
                logger.info(f"[确认等待] {tool_name} pending_id={pending_id}")
                return {
                    "content": preview,
                    "need_confirm": True,
                    "pending_id": pending_id,
                }

            # 安全操作 或 已确认 → 直接执行
            logger.info(f"执行工具: {tool_name}({args})")
            import asyncio
            result = handler(args, user_id)
            if asyncio.iscoroutine(result):
                result = await result
            logger.info(f"工具结果: {str(result)[:200]}")
            return {"content": str(result), "need_confirm": False}

        except Exception as e:
            logger.error(f"工具执行异常 [{tool_name}]: {e}")
            return {"content": f"工具执行出错: {e}", "need_confirm": False}

    # ==================== 确认机制 ====================

    def _save_pending_action(self, user_id: str, tool_name: str, args: dict) -> str:
        """将待确认操作存入 SQLite，返回 pending_id"""
        import uuid
        pending_id = str(uuid.uuid4())[:8]
        conn = self._get_pending_db()
        try:
            conn.execute(
                "INSERT INTO pending_actions (id, user_id, tool_name, args_json, created_at) VALUES (?, ?, ?, ?, ?)",
                (pending_id, user_id, tool_name, json.dumps(args, ensure_ascii=False), datetime.now().isoformat()),
            )
            conn.commit()
        finally:
            conn.close()
        return pending_id

    def get_pending_action(self, pending_id: str) -> Optional[dict]:
        """根据 ID 获取待确认操作"""
        conn = self._get_pending_db()
        try:
            row = conn.execute(
                "SELECT id, user_id, tool_name, args_json FROM pending_actions WHERE id = ?",
                (pending_id,),
            ).fetchone()
            if row:
                return {"id": row[0], "user_id": row[1], "tool_name": row[2], "args": json.loads(row[3])}
            return None
        finally:
            conn.close()

    def remove_pending_action(self, pending_id: str):
        """删除已处理的待确认操作"""
        conn = self._get_pending_db()
        try:
            conn.execute("DELETE FROM pending_actions WHERE id = ?", (pending_id,))
            conn.commit()
        finally:
            conn.close()

    def _get_pending_db(self) -> sqlite3:
        """获取 pending_actions 表的连接"""
        db_path = os.path.join(os.path.dirname(MEMORY_DB_PATH), "memory.db")
        conn = sqlite3.connect(str(db_path), timeout=10)
        conn.execute("""CREATE TABLE IF NOT EXISTS pending_actions (
            id TEXT PRIMARY KEY, user_id TEXT NOT NULL, tool_name TEXT NOT NULL,
            args_json TEXT NOT NULL, created_at TEXT NOT NULL
        )""")
        conn.commit()
        return conn

    def _preview_dangerous_action(self, tool_name: str, args: dict) -> str:
        """生成危险操作的预览描述"""
        if tool_name == "delete_file":
            path = args.get("path", args.get("file", ""))
            return f"⚠️ **即将删除**: `{path}`\n\n此操作不可撤销，确认要删除吗？"
        elif tool_name == "run_command":
            cmd = args.get("command", args.get("cmd", ""))
            return f"⚠️ **即将执行命令**: `{cmd}`\n\n确认要执行吗？"
        return f"⚠️ 即将执行 **{tool_name}**，确认吗？"

    # ==================== 天气查询 ====================

    # IP 定位缓存（类级别，启动时预加载）
    _cached_location = None

    @classmethod
    def preload_location(cls):
        """启动时预加载IP定位，避免首次查询等待"""
        try:
            cls._get_location_by_ip()
        except Exception:
            pass

    @classmethod
    def _get_location_by_ip(cls) -> dict:
        """通过 IP 获取当前城市和坐标，优先国内源"""
        if cls._cached_location:
            return cls._cached_location

        # 源1: api.ip.sb（国内快）
        try:
            r = httpx.get("https://api.ip.sb/geoip", timeout=3)
            d = r.json()
            city = d.get("city", "")
            lat = d.get("latitude", 0)
            lon = d.get("longitude", 0)
            if city and lat:
                result = {"city": city, "lat": lat, "lon": lon}
                cls._cached_location = result
                logger.info(f"IP定位: {city} ({lat}, {lon})")
                return result
        except Exception:
            pass

        # 源2: ip-api.com
        try:
            r = httpx.get("http://ip-api.com/json/?lang=zh-CN&fields=city,lat,lon", timeout=3)
            d = r.json()
            city = d.get("city", "")
            lat = d.get("lat", 0)
            lon = d.get("lon", 0)
            if city and lat:
                result = {"city": city, "lat": lat, "lon": lon}
                cls._cached_location = result
                logger.info(f"IP定位: {city} ({lat}, {lon})")
                return result
        except Exception:
            pass

        # 兜底
        fallback = {"city": "北京", "lat": 39.9, "lon": 116.4}
        cls._cached_location = fallback
        logger.warning("IP定位失败，使用默认城市：北京")
        return fallback

    @staticmethod
    def _geocode(city: str) -> tuple:
        """用 open-meteo 免费地理编码 API 把城市名转经纬度（全球城市）"""
        # 先从本地缓存查
        if city in CITY_COORDS:
            return CITY_COORDS[city]
        try:
            r = httpx.get(
                "https://geocoding-api.open-meteo.com/v1/search",
                params={"name": city, "count": 1, "language": "zh", "format": "json"},
                timeout=5,
            )
            data = r.json()
            results = data.get("results", [])
            if results:
                return (results[0]["latitude"], results[0]["longitude"])
        except Exception:
            pass
        return None

    async def _tool_get_weather(self, args: dict, user_id: str = "") -> str:
        city = args.get("city", args.get("location", ""))
        date_str = args.get("date", "")

        # 没指定城市 → 用户记忆 > IP定位
        if not city:
            # 查用户记忆
            if self.memory_service and user_id:
                try:
                    mems = await self.memory_service.search_memories(user_id, "所在城市")
                    if mems:
                        city = mems[0].replace("[所在城市]", "").strip()
                except Exception:
                    pass
            # IP 兜底
            if not city:
                city = self._get_location_by_ip()["city"]

        coords = self._geocode(city)
        if not coords:
            return f"未找到城市「{city}」。如需设为默认城市：对我说'记住我在XX'"
        lat, lon = coords
        try:
            resp = httpx.get(
                "https://api.open-meteo.com/v1/forecast",
                params={
                    "latitude": lat, "longitude": lon,
                    "daily": "weather_code,temperature_2m_max,temperature_2m_min,precipitation_sum,wind_speed_10m_max",
                    "timezone": "Asia/Shanghai", "forecast_days": 7,
                },
                timeout=8,
            )
            data = resp.json()
        except Exception as e:
            return f"天气数据获取失败: {e}"

        daily = data.get("daily", {})
        dates = daily.get("time", [])
        if not dates:
            return "天气数据暂时不可用"

        t_max = daily.get("temperature_2m_max", [])
        t_min = daily.get("temperature_2m_min", [])
        weather_codes = daily.get("weather_code", [])
        wind = daily.get("wind_speed_10m_max", [])
        rain = daily.get("precipitation_sum", [])

        lines = [f"**{city}** 天气预报\n"]
        today = datetime.now().strftime("%Y-%m-%d")
        tomorrow = (datetime.now() + timedelta(days=1)).strftime("%Y-%m-%d")

        for i, d in enumerate(dates):
            if date_str and date_str != d:
                continue
            w_desc = WEATHER_CODES.get(weather_codes[i], "未知") if i < len(weather_codes) else "?"
            high = t_max[i] if i < len(t_max) else "?"
            low = t_min[i] if i < len(t_min) else "?"
            w = wind[i] if i < len(wind) else "?"
            r = rain[i] if i < len(rain) else "?"

            day_label = d
            if d == today:
                day_label = f"{d}（今天）"
            elif d == tomorrow:
                day_label = f"{d}（明天）"

            lines.append(f"**{day_label}**: {w_desc}，{low}°C ~ {high}°C，风速{w}km/h，降水{r}mm")
            if date_str:
                break

        return "\n".join(lines)

    # ==================== 计算器 ====================

    def _tool_calculator(self, args: dict, user_id: str = "") -> str:
        expression = args.get("expression", "")
        if not expression:
            return "请提供要计算的表达式"
        try:
            result = _safe_eval(expression)
            return f"**{expression} = {result}**"
        except ValueError as e:
            return f"计算失败: {e}"
        except Exception as e:
            return f"计算出错: {e}"

    # ==================== 当前时间 ====================

    def _tool_get_current_time(self, args: dict, user_id: str = "") -> str:
        tz_name = args.get("timezone", "Asia/Shanghai")
        tz_map = {"Asia/Shanghai": 8, "Asia/Tokyo": 9, "America/New_York": -4, "Europe/London": 1, "UTC": 0}
        offset = tz_map.get(tz_name, 8)
        now = datetime.now(timezone(timedelta(hours=offset)))
        weekdays = ["星期一", "星期二", "星期三", "星期四", "星期五", "星期六", "星期日"]
        return f"**{now.strftime('%Y年%m月%d日')} {weekdays[now.weekday()]} {now.strftime('%H:%M:%S.%f')}**（{tz_name}）"

    # ==================== 记忆 ====================

    async def _tool_save_memory(self, args: dict, user_id: str = "") -> str:
        if not self.memory_service:
            return "记忆服务不可用"
        key = args.get("key", "")
        content = args.get("content", "")
        if not key or not content:
            return "请提供记忆标签和内容"
        await self.memory_service.save_memory(user_id, key, content)
        return f"已记住: **{key}** → {content}"

    async def _tool_search_memory(self, args: dict, user_id: str = "") -> str:
        if not self.memory_service:
            return "记忆服务不可用"
        keyword = args.get("keyword", args.get("query", ""))
        results = await self.memory_service.search_memories_semantic(user_id, keyword)
        if not results:
            return "未找到相关记忆"
        return "\n".join(results)

    # ==================== 文档检索 ====================

    def _tool_search_documents(self, args: dict, user_id: str = "") -> str:
        if not self.rag_service or not self.rag_service.is_ready:
            return "文档检索服务不可用。请先上传文档。"
        query = args.get("query", "")
        results = self.rag_service.search(user_id, query)
        if not results:
            return "未找到相关文档内容。"
        return "\n\n".join(f"[片段{i+1}] {r}" for i, r in enumerate(results))

    # ==================== 文件操作 ====================

    def _tool_list_directory(self, args: dict, user_id: str = "") -> str:
        path = args.get("path", args.get("dir", ""))
        if not path:
            return "请提供目录路径，例如: path='C:/Users/李琳/Desktop'"

        real = self._resolve_path(path)
        if not real:
            return f"路径不存在: {path}"

        if not os.path.isdir(real):
            return f"不是目录: {path}"

        try:
            items = os.listdir(real)
            if not items:
                return f"目录为空: {path}"

            lines = [f"**{path}** ({len(items)}项)\n"]
            for item in sorted(items)[:80]:
                full = os.path.join(real, item)
                try:
                    if os.path.isdir(full):
                        lines.append(f"  📂 {item}/")
                    else:
                        size = os.path.getsize(full)
                        size_str = f"{size/1024:.1f}KB" if size > 1024 else f"{size}B"
                        lines.append(f"  📄 {item} ({size_str})")
                except OSError:
                    lines.append(f"  🔒 {item}")

            if len(items) > 80:
                lines.append(f"\n  ... 还有{len(items)-80}项")
            return "\n".join(lines)
        except PermissionError:
            return f"没有权限: {path}"
        except Exception as e:
            return f"读取失败: {e}"

    def _tool_read_file(self, args: dict, user_id: str = "") -> str:
        path = args.get("path", args.get("file", ""))
        if not path:
            return "请提供文件路径"

        real = self._resolve_path(path)
        if not real:
            return f"文件不存在: {path}"

        if os.path.isdir(real):
            return f"这是目录不是文件: {path}"

        if os.path.getsize(real) > 5 * 1024 * 1024:
            return f"文件太大（>5MB）: {path}"

        ext = Path(real).suffix.lower()

        # PDF 文件
        if ext == ".pdf":
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(real)
                pages = []
                for i, page in enumerate(reader.pages[:50]):
                    text = page.extract_text()
                    if text:
                        pages.append(f"--- 第{i+1}页 ---\n{text}")
                content = "\n\n".join(pages) if pages else "PDF 文件无法提取文本"
                return f"**{path}** ({len(reader.pages)}页)\n\n{content[:20000]}"
            except Exception as e:
                return f"PDF 读取失败: {e}"

        # Word 文件
        if ext in (".docx", ".doc"):
            try:
                from docx import Document
                doc = Document(real)
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                content = "\n".join(paragraphs)
                # 也读取表格
                tables_text = []
                for table in doc.tables:
                    for row in table.rows:
                        cells = [cell.text for cell in row.cells]
                        tables_text.append(" | ".join(cells))
                if tables_text:
                    content += "\n\n--- 表格 ---\n" + "\n".join(tables_text)
                return f"**{path}**\n\n{content[:20000]}"
            except Exception as e:
                return f"Word 文件读取失败: {e}"

        # PPTX 文件
        if ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(real)
                slides_text = []
                for i, slide in enumerate(prs.slides):
                    lines = [f"--- 第{i+1}页 ---"]
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                text = para.text.strip()
                                if text:
                                    lines.append(text)
                        if shape.has_table:
                            table = shape.table
                            for row in table.rows:
                                cells = [cell.text.strip() for cell in row.cells]
                                lines.append(" | ".join(cells))
                    slides_text.append("\n".join(lines))
                content = "\n\n".join(slides_text)
                return f"**{path}** ({len(prs.slides)}页)\n\n{content[:20000]}"
            except Exception as e:
                return f"PPTX 读取失败: {e}"

        # 普通文本文件
        try:
            with open(real, "r", encoding="utf-8", errors="replace") as f:
                content = f.read(20000)
            return f"**{path}**\n\n{content}"
        except Exception as e:
            return f"读取失败: {e}"

    def _tool_write_file(self, args: dict, user_id: str = "") -> str:
        path = args.get("path", args.get("file", ""))
        content = args.get("content", "")
        if not path:
            return "请提供文件路径和内容"

        # 解析路径（允许新文件）
        target = os.path.expanduser(path)
        target = os.path.expandvars(target)
        target = os.path.normpath(target)

        if self._is_protected(target):
            return f"不允许写入系统目录: {path}"

        try:
            parent = os.path.dirname(target)
            if parent:
                os.makedirs(parent, exist_ok=True)
            with open(target, "w", encoding="utf-8") as f:
                f.write(content)
            return f"✅ 已创建/写入文件: **{path}** ({len(content)}字符)"
        except Exception as e:
            return f"写入失败: {e}"

    def _tool_delete_file(self, args: dict, user_id: str = "") -> str:
        path = args.get("path", args.get("file", ""))
        if not path:
            return "请提供要删除的文件路径"

        real = self._resolve_path(path)
        if not real:
            return f"文件不存在: {path}"

        if self._is_protected(real):
            return f"⛔ 禁止删除系统文件: {path}"

        if not os.path.exists(real):
            return f"不存在: {path}"

        try:
            if os.path.isdir(real):
                shutil.rmtree(real)
                return f"✅ 已删除文件夹: **{path}**"
            else:
                os.remove(real)
                return f"✅ 已删除文件: **{path}**"
        except PermissionError:
            return f"没有删除权限: {path}"
        except Exception as e:
            return f"删除失败: {e}"

    # ==================== 打开文件/应用 ====================

    def _tool_open_file(self, args: dict, user_id: str = "") -> str:
        """
        用系统默认程序打开文件，或启动应用程序
        Windows 用 os.startfile，Linux/Mac 用 subprocess
        """
        path = args.get("path", args.get("file", args.get("app", args.get("application", ""))))
        if not path:
            return "请提供文件路径或应用程序名称"

        # 尝试直接路径
        real = self._resolve_path(path)

        # 如果不是直接路径，尝试作为应用程序名搜索
        if not real:
            app_path = shutil.which(path)
            if not app_path:
                # Windows 常见应用搜索
                app_path = self._find_windows_app(path)
            if app_path:
                real = app_path

        if not real:
            return f"找不到文件或应用: {path}"

        try:
            if os.name == "nt":  # Windows
                os.startfile(real)
            else:  # Linux/Mac
                subprocess.Popen(["xdg-open" if os.name != "darwin" else "open", real])
            return f"✅ 已打开: **{path}**"
        except Exception as e:
            return f"打开失败: {e}"

    def _find_windows_app(self, name: str) -> str:
        """在 Windows 常见路径中搜索应用程序（增强版）"""
        name_lower = name.lower().replace(".exe", "").strip()

        # 中文名→英文目录名别名
        _NAME_ALIASES = {
            "qq音乐": ["qqmusic", "qq music"],
            "网易云音乐": ["cloudmusic", "cloud music", "wangyiyun"],
            "网易云": ["cloudmusic", "cloud music", "wangyiyun"],
            "腾讯会议": ["wemeet", "tencent meeting"],
            "钉钉": ["dingtalk", "dingding", "ding talk"],
            "微信": ["wechat"],
            "计算器": ["calc", "calculator"],
            "记事本": ["notepad"],
            "命令提示符": ["cmd", "terminal"],
        }
        # 扩展搜索词：原名 + 别名
        search_names = {name_lower}
        for alias_key, alias_list in _NAME_ALIASES.items():
            if name_lower == alias_key or alias_key in name_lower or name_lower in alias_key:
                search_names.update(a.lower() for a in alias_list)
        # 也尝试拼音匹配
        search_names.add(name_lower)

        def _matches_dir(dirname: str) -> bool:
            """检查目录名是否匹配任一搜索词"""
            d = dirname.lower()
            for sn in search_names:
                if sn in d or d in sn:
                    return True
            return False

        # ---- 1. 常见应用硬编码映射（快速路径）----
        app_map = {
            # 浏览器
            "chrome": r"C:\Program Files\Google\Chrome\Application\chrome.exe",
            "firefox": r"C:\Program Files\Mozilla Firefox\firefox.exe",
            "edge": r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
            # 腾讯系
            "微信": r"C:\Program Files\Tencent\WeChat\WeChat.exe",
            "wechat": r"C:\Program Files\Tencent\WeChat\WeChat.exe",
            "qq": r"C:\Program Files\Tencent\QQ\Bin\QQ.exe",
            "qq音乐": r"C:\Program Files (x86)\Tencent\QQMusic\QQMusic.exe",
            "qqmusic": r"C:\Program Files (x86)\Tencent\QQMusic\QQMusic.exe",
            # 网易系
            "网易云音乐": r"C:\Program Files (x86)\NetEase\CloudMusic\cloudmusic.exe",
            "网易云": r"C:\Program Files (x86)\NetEase\CloudMusic\cloudmusic.exe",
            "netsease": r"C:\Program Files (x86)\NetEase\CloudMusic\cloudmusic.exe",
            # 开发工具
            "vscode": r"C:\Users\李琳\AppData\Local\Programs\Microsoft VS Code\new_Code.exe",
            "code": r"C:\Users\李琳\AppData\Local\Programs\Microsoft VS Code\new_Code.exe",
            # 系统
            "notepad": "notepad.exe",
            "记事本": "notepad.exe",
            "calc": "calc.exe",
            "计算器": "calc.exe",
            "explorer": "explorer.exe",
            "cmd": "cmd.exe",
            "命令提示符": "cmd.exe",
            "任务管理器": "taskmgr.exe",
            "taskmgr": "taskmgr.exe",
            # Office
            "word": r"C:\Program Files\Microsoft Office\root\Office16\WINWORD.EXE",
            "excel": r"C:\Program Files\Microsoft Office\root\Office16\EXCEL.EXE",
            "ppt": r"C:\Program Files\Microsoft Office\root\Office16\POWERPNT.EXE",
            "powerpoint": r"C:\Program Files\Microsoft Office\root\Office16\POWERPNT.EXE",
            "wps": r"C:\Users\李琳\AppData\Local\Kingsoft\WPS Office\wps.exe",
            # 其他
            "typora": r"C:\Program Files\Typora\Typora.exe",
            "钉钉": r"C:\Program Files (x86)\DingDing\Dingtalk.exe",
            "dingtalk": r"C:\Program Files (x86)\DingDing\Dingtalk.exe",
            "腾讯会议": r"C:\Program Files (x86)\Tencent\WeMeet\wemeetapp.exe",
            "wemeet": r"C:\Program Files (x86)\Tencent\WeMeet\wemeetapp.exe",
            "tim": r"C:\Program Files (x86)\Tencent\TIM\Bin\TIM.exe",
        }

        # 精确匹配
        if name_lower in app_map:
            p = app_map[name_lower]
            if os.path.exists(p) or shutil.which(p):
                return p if os.path.exists(p) else shutil.which(p)

        # 模糊匹配
        for key, p in app_map.items():
            if name_lower in key or key in name_lower:
                if os.path.exists(p) or shutil.which(p):
                    return p if os.path.exists(p) else shutil.which(p)

        # ---- 2. PATH 搜索 ----
        for ext in [".exe", ".lnk", ".url"]:
            result = shutil.which(name + ext) or shutil.which(name)
            if result and os.path.exists(result):
                return result

        # ---- 3. 桌面快捷方式 ----
        desktop = os.path.expanduser("~/Desktop")
        try:
            for f in os.listdir(desktop):
                if _matches_dir(f):
                    p = os.path.join(desktop, f)
                    # .lnk 文件直接返回（os.startfile 可以打开）
                    if os.path.exists(p):
                        return p
        except Exception:
            pass

        # ---- 4. 用户开始菜单 ----
        for start_menu in [
            os.path.expanduser(r"~\AppData\Roaming\Microsoft\Windows\Start Menu\Programs"),
            r"C:\ProgramData\Microsoft\Windows\Start Menu\Programs",
        ]:
            try:
                for root, dirs, files in os.walk(start_menu):
                    for f in files:
                        if _matches_dir(f):
                            return os.path.join(root, f)
            except Exception:
                continue

        # ---- 5. Program Files 搜索（含腾讯专属子目录）----
        for base in [r"C:\Program Files", r"C:\Program Files (x86)"]:
            try:
                for entry in os.listdir(base):
                    entry_path = os.path.join(base, entry)
                    if not os.path.isdir(entry_path):
                        continue
                    # 直接匹配
                    if _matches_dir(entry):
                        for f in os.listdir(entry_path):
                            if f.lower().endswith(".exe") and _matches_dir(os.path.splitext(f)[0]):
                                return os.path.join(entry_path, f)
                        for f in os.listdir(entry_path):
                            if f.lower().endswith(".exe"):
                                return os.path.join(entry_path, f)
                    # 腾讯/网易等大厂子目录深度搜索
                    if entry.lower() in ("tencent", "netease", "kingsoft", "wps office",
                                          "google", "mozilla", "microsoft", "dingding",
                                          "alibaba", "baidu", "360", "sogou", "xunlei", "bilibili"):
                        for subdir in os.listdir(entry_path):
                            subdir_path = os.path.join(entry_path, subdir)
                            if not os.path.isdir(subdir_path):
                                continue
                            if _matches_dir(subdir):
                                for root, dirs, files in os.walk(subdir_path):
                                    for f in files:
                                        if f.lower().endswith(".exe") and _matches_dir(os.path.splitext(f)[0]):
                                            return os.path.join(root, f)
                                # fallback: return first exe found
                                for root, dirs, files in os.walk(subdir_path):
                                    for f in files:
                                        if f.lower().endswith(".exe"):
                                            return os.path.join(root, f)
            except Exception:
                continue

        # ---- 6. Local AppData 搜索 ----
        local_appdata = os.path.expanduser(r"~\AppData\Local")
        try:
            for entry in os.listdir(local_appdata):
                if _matches_dir(entry):
                    entry_path = os.path.join(local_appdata, entry)
                    if os.path.isdir(entry_path):
                        for root, dirs, files in os.walk(entry_path):
                            for f in files:
                                if f.lower().endswith(".exe") and _matches_dir(os.path.splitext(f)[0]):
                                    return os.path.join(root, f)
        except Exception:
            pass

        return ""

    # ==================== 执行命令 ====================

    def _tool_run_command(self, args: dict, user_id: str = "") -> str:
        """执行系统命令（有安全限制）"""
        command = args.get("command", args.get("cmd", ""))
        if not command:
            return "请提供要执行的命令"

        # 安全黑名单
        blocked = ["format", "del /", "rmdir /s", "rd /s", "shutdown", "taskkill /f",
                    "reg delete", "reg add", "net user", "net localgroup",
                    "takeown /f c:", "cipher /w", "diskpart"]
        cmd_lower = command.lower()
        for b in blocked:
            if b in cmd_lower:
                return f"⛔ 禁止执行危险命令: {command}"

        try:
            result = subprocess.run(
                command,
                shell=True,
                capture_output=True,
                text=True,
                timeout=30,
                encoding="utf-8",
                errors="replace",
            )
            output = result.stdout.strip()
            error = result.stderr.strip()

            if result.returncode == 0:
                msg = f"✅ 命令执行成功"
                if output:
                    msg += f":\n{output[:3000]}"
                return msg
            else:
                msg = f"⚠️ 命令返回码 {result.returncode}"
                if error:
                    msg += f":\n{error[:1000]}"
                return msg
        except subprocess.TimeoutExpired:
            return "⏱️ 命令执行超时（30秒限制）"
        except Exception as e:
            return f"执行失败: {e}"

    # ==================== 网页搜索 ====================

    def _tool_web_search(self, args: dict, user_id: str = "") -> str:
        query = args.get("query", args.get("q", args.get("keyword", "")))
        if not query:
            return "请提供搜索关键词"

        import concurrent.futures
        from urllib.parse import quote

        def _search_bing_cn():
            """Bing 中国版，国内可直接访问"""
            try:
                resp = httpx.get(
                    f"https://cn.bing.com/search?q={quote(query)}&count=5",
                    headers={
                        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
                        "Accept-Language": "zh-CN,zh;q=0.9",
                    },
                    timeout=6,
                    follow_redirects=True,
                )
                blocks = re.findall(r'<li class="b_algo"[^>]*>(.*?)</li>', resp.text, re.DOTALL)
                results = []
                for block in blocks[:5]:
                    title_m = re.search(r'<h2[^>]*><a[^>]*>(.*?)</a>', block, re.DOTALL)
                    snippet_m = re.search(r'(?:<p class="b_lineclamp[^"]*"[^>]*>|<p[^>]*>)(.*?)</p>', block, re.DOTALL)
                    if title_m:
                        t = re.sub(r'<[^>]+>', '', title_m.group(1)).strip()
                        s = re.sub(r'<[^>]+>', '', snippet_m.group(1)).strip() if snippet_m else ""
                        results.append(f"{len(results)+1}. {t}\n   {s[:200]}")
                return results if results else None
            except Exception:
                return None

        def _search_sogou():
            """搜狗搜索，国内源"""
            try:
                resp = httpx.get(
                    f"https://www.sogou.com/web?query={quote(query)}",
                    headers={
                        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36",
                    },
                    timeout=6,
                )
                # 搜狗结果在 class="vrwrap" 或 class="rb" 中
                blocks = re.findall(r'<div[^>]*class="(?:vrwrap|rb)"[^>]*>(.*?)</div>\s*</div>', resp.text, re.DOTALL)
                results = []
                for block in blocks[:5]:
                    title_m = re.search(r'<a[^>]*id="[^"]*"[^>]*>(.*?)</a>', block, re.DOTALL)
                    snippet_m = re.search(r'<p[^>]*class="(?:star-wiki|str_info|str-text)[^"]*"[^>]*>(.*?)</p>', block, re.DOTALL)
                    if title_m:
                        t = re.sub(r'<[^>]+>', '', title_m.group(1)).strip()
                        s = re.sub(r'<[^>]+>', '', snippet_m.group(1)).strip() if snippet_m else ""
                        if t:
                            results.append(f"{len(results)+1}. {t}\n   {s[:200]}")
                return results if results else None
            except Exception:
                return None

        # 双源竞速（总超时 8 秒）
        try:
            with concurrent.futures.ThreadPoolExecutor(max_workers=2) as ex:
                futures = {
                    ex.submit(_search_bing_cn): "Bing",
                    ex.submit(_search_sogou): "搜狗",
                }
                for fut in concurrent.futures.as_completed(futures, timeout=8):
                    try:
                        results = fut.result()
                        if results:
                            return f"搜索 **{query}**\n\n" + "\n\n".join(results)
                    except Exception:
                        continue
        except concurrent.futures.TimeoutError:
            pass

        return f"搜索「{query}」超时，请换个简短关键词重试。"

    # ==================== 图片生成 ====================

    def _tool_generate_image(self, args: dict, user_id: str = "") -> str:
        """生成图片：硅基流动优先，智谱 CogView-4 备选"""
        prompt = args.get("prompt", "")
        size = args.get("size", "1024x1024")
        if not prompt:
            return "请提供图片描述"

        # 优先使用硅基流动
        if SILICONFLOW_API_KEY:
            try:
                resp = httpx.post(
                    f"{SILICONFLOW_BASE_URL}/images/generations",
                    headers={
                        "Authorization": f"Bearer {SILICONFLOW_API_KEY}",
                        "Content-Type": "application/json",
                    },
                    json={
                        "model": SILICONFLOW_IMAGE_MODEL,
                        "prompt": prompt,
                        "image_size": size,
                        "num_inference_steps": 8,
                    },
                    timeout=60,
                )
                data = resp.json()

                if resp.status_code == 200 and "images" in data:
                    images = data["images"]
                    if images:
                        url = images[0].get("url", "")
                        if url:
                            return f"✅ 图片生成成功！\n\n**描述**: {prompt}\n**尺寸**: {size}\n\n![生成的图片]({url})\n\n图片链接: {url}"
                error_msg = data.get("message", data.get("error", {}).get("message", str(data)))
                logger.warning(f"[图片生成] 硅基流动失败: {error_msg}，尝试智谱备选")
            except Exception as e:
                logger.warning(f"[图片生成] 硅基流动异常: {e}，尝试智谱备选")

        # 备选：智谱 CogView-4
        if ZHIPU_API_KEY:
            try:
                resp = httpx.post(
                    f"{ZHIPU_BASE_URL}/images/generations",
                    headers={
                        "Authorization": f"Bearer {ZHIPU_API_KEY}",
                        "Content-Type": "application/json",
                    },
                    json={
                        "model": "cogview-4",
                        "prompt": prompt,
                        "size": size,
                    },
                    timeout=120,
                )
                data = resp.json()

                if resp.status_code == 200:
                    images = data.get("data", [])
                    if images:
                        url = images[0].get("url", "")
                        if url:
                            return f"✅ 图片生成成功！\n\n**描述**: {prompt}\n**尺寸**: {size}\n\n![生成的图片]({url})\n\n图片链接: {url}"
                error_msg = data.get("message", data.get("error", {}).get("message", str(data)))
                return f"图片生成失败: {error_msg}"
            except Exception as e:
                return f"图片生成出错: {e}"

        return "图片生成服务不可用（未配置硅基流动或智谱 API Key）"

    # ==================== 语音识别 ====================

    def transcribe_audio(self, audio_data: bytes) -> str:
        """使用硅基流动语音识别"""
        if not SILICONFLOW_API_KEY:
            return "语音识别服务未配置"

        try:
            import tempfile
            with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as f:
                f.write(audio_data)
                tmp_path = f.name

            with open(tmp_path, "rb") as f:
                resp = httpx.post(
                    f"{SILICONFLOW_BASE_URL}/audio/transcriptions",
                    headers={"Authorization": f"Bearer {SILICONFLOW_API_KEY}"},
                    files={"file": ("audio.wav", f, "audio/wav")},
                    data={"model": "FunAudioLLM/SenseVoiceSmall"},
                    timeout=30,
                )

            os.unlink(tmp_path)
            data = resp.json()
            return data.get("text", "")
        except Exception as e:
            return f"语音识别失败: {e}"

    # ==================== UI 自动化 ====================

    def _tool_app_get_controls(self, args: dict, user_id: str = "") -> str:
        """获取窗口 UI 控件列表（pywinauto）"""
        window_title = args.get("window_title", "")
        max_depth = int(args.get("max_depth", 3))

        if not window_title:
            return "请提供窗口标题关键词，如：'记事本'、'QQ音乐'"

        try:
            from pywinauto import Desktop
        except ImportError:
            return "pywinauto 未安装，请执行 pip install pywinauto"

        import time

        # 重试搜索：窗口刚打开时 UIA 可能还没注册
        dlg = None
        found_titles = []  # 记录找到的所有窗口标题，帮助诊断
        for attempt in range(6):  # 最多重试 6 次（共 ~3 秒）
            for backend in ["uia", "win32"]:
                try:
                    desktop = Desktop(backend=backend)
                    best_match = None  # 优先选有意义的标题
                    for w in desktop.windows():
                        try:
                            wt = w.window_text()
                            if not wt or len(wt) == 0:
                                continue
                            if attempt == 0:
                                found_titles.append(wt[:80])
                            # 跳过 GDI+/IME 等内部窗口
                            if wt.startswith("GDI+") or wt.startswith("IME") or wt == "Default IME":
                                continue
                            wt_lower = wt.lower()
                            title_lower = window_title.lower()
                            # 精确匹配最佳
                            if title_lower in wt_lower:
                                dlg = w
                                break
                            # 反向匹配
                            if wt_lower in title_lower:
                                best_match = best_match or w
                            # 特殊别名
                            if title_lower in ("记事本", "notepad") and \
                               "notepad" in wt_lower:
                                dlg = w
                                break
                            if title_lower in ("计算器", "calc", "calculator") and \
                               ("calc" in wt_lower or "calculator" in wt_lower):
                                dlg = w
                                break
                        except Exception:
                            continue
                    if dlg:
                        break
                    if best_match:
                        dlg = best_match
                        break
                except Exception:
                    continue
            if dlg:
                break
            time.sleep(0.5)  # 等窗口就绪

        if not dlg:
            # 没找到时，列出实际存在的窗口标题帮助用户
            unique_titles = list(set(found_titles))[:10]
            titles_str = "\n  ".join(t for t in unique_titles if t.strip())
            return (
                f"未找到包含'{window_title}'的窗口。当前可见窗口：\n"
                f"  {titles_str}\n"
                f"提示：请用上面显示的精确标题重试，或手动点击目标窗口使其置前。"
            )

        # 获取控件树（后续逻辑不变）
        try:
            dlg.wait("visible", timeout=3)
        except Exception:
            pass
            controls = dlg.descendants()
            if not controls:
                return f"窗口 '{window_title}' 没有可访问的UI控件"

            # 过滤+格式化输出
            meaningful_types = {
                "Button", "Edit", "ComboBox", "ListItem", "ListBox",
                "MenuItem", "CheckBox", "RadioButton", "TabItem",
                "TreeItem", "Hyperlink", "Text", "Slider", "Spinner",
            }

            lines = [f"**窗口 '{window_title}' 的UI控件** ({len(controls)}个)\n"]
            count = 0
            limit = 50

            for c in controls:
                try:
                    c_type = c.element_info.control_type if hasattr(c, 'element_info') else type(c).__name__
                    c_name = c.window_text()[:60] if c.window_text() else "(无文本)"
                    c_auto = c.element_info.automation_id if hasattr(c, 'element_info') and c.element_info.automation_id else ""

                    # 只输出有意义的控件
                    if c_type not in meaningful_types and c_type != "Pane":
                        continue

                    # 跳过太深的嵌套容器
                    if c_type == "Pane" and not c_name.strip():
                        continue

                    auto_str = f" [{c_auto}]" if c_auto else ""
                    lines.append(f"  {count+1}. [{c_type}] {c_name}{auto_str}")
                    count += 1
                    if count >= limit:
                        lines.append(f"\n  ... 还有 {len(controls) - count} 个控件未展示")
                        break
                except Exception:
                    continue

            if count == 0:
                return f"窗口 '{window_title}' 没有可操作的UI控件（可能被遮挡或最小化）"

            return "\n".join(lines)

        except Exception as e:
            return f"获取控件列表失败: {e}"

    def _tool_app_action(self, args: dict, user_id: str = "") -> str:
        """对窗口 UI 控件执行操作（pywinauto）"""
        window_title = args.get("window_title", "")
        target = args.get("target", "")
        action = args.get("action", "click")
        value = args.get("value", "")

        if not window_title:
            return "请提供目标窗口标题"
        if action not in [
            "click", "double_click", "right_click", "type_text",
            "get_text", "select", "scroll_down", "scroll_up",
            "maximize", "minimize", "wait_for",
        ]:
            return f"不支持的操作类型: {action}"

        try:
            from pywinauto import Desktop
            from pywinauto.keyboard import send_keys
        except ImportError:
            return "pywinauto 未安装，请执行 pip install pywinauto"

        try:
            import time

            # 查找窗口（多重匹配 + 重试）
            dlg = None
            for attempt in range(6):
                for backend in ["uia", "win32"]:
                    try:
                        desktop = Desktop(backend=backend)
                        for w in desktop.windows():
                            try:
                                wt = w.window_text()
                                if not wt:
                                    continue
                                wt_lower = wt.lower()
                                title_lower = window_title.lower()
                                if title_lower in wt_lower or wt_lower in title_lower:
                                    dlg = w
                                    break
                                # 特殊处理常见中文名
                                if title_lower in ("记事本", "notepad") and \
                                   ("记事本" in wt or "notepad" in wt_lower):
                                    dlg = w
                                    break
                                if title_lower in ("计算器", "calc") and \
                                   ("计算器" in wt or "calc" in wt_lower or "calculator" in wt_lower):
                                    dlg = w
                                    break
                            except Exception:
                                continue
                        if dlg:
                            break
                    except Exception:
                        continue
                if dlg:
                    break
                time.sleep(0.5)

            if not dlg:
                return f"未找到包含'{window_title}'的窗口。请确认软件已打开，并尝试用 app_get_controls 查看实际窗口标题。"

            # 确保窗口可见
            try:
                dlg.set_focus()
            except Exception:
                pass

            # 如果无 target，直接操作窗口本身
            if not target:
                if action == "maximize":
                    dlg.maximize()
                    return f"✅ 窗口 '{window_title}' 已最大化"
                elif action == "minimize":
                    dlg.minimize()
                    return f"✅ 窗口 '{window_title}' 已最小化"
                elif action == "click":
                    dlg.click()
                    return f"✅ 已点击窗口 '{window_title}'"
                elif action == "type_text":
                    return self._do_type_text(value, window_title)
                else:
                    return f"操作 '{action}' 需要指定 target（目标控件）"

            # ---- 查找目标控件 ----
            ctrl = None
            try:
                ctrl = dlg.child_window(title=target)
            except Exception:
                pass

            if not ctrl:
                for c in dlg.descendants():
                    try:
                        c_name = c.window_text() if c.window_text() else ""
                        c_auto = c.element_info.automation_id if hasattr(c, 'element_info') and c.element_info.automation_id else ""
                        if target.lower() in c_name.lower() or target == c_auto:
                            ctrl = c
                            break
                    except Exception:
                        continue

            if not ctrl:
                return (
                    f"未找到控件 '{target}'。\n"
                    f"建议改用 app_keyboard 或 app_clipboard 方案。"
                )

            # ---- 执行操作 ----
            if action == "click":
                ctrl.click()
                return f"✅ 已点击 '{target}'"
            elif action == "double_click":
                ctrl.double_click()
                return f"✅ 已双击 '{target}'"
            elif action == "right_click":
                ctrl.right_click()
                return f"✅ 已右键点击 '{target}'"
            elif action == "type_text":
                if not value:
                    return "type_text 操作需要提供 value 参数（要输入的文字）"
                try:
                    ctrl.type_keys(value, pause=0.05)
                    return f"✅ 已在 '{target}' 中输入: {value[:50]}"
                except Exception:
                    return self._do_type_text(value, window_title)
            elif action == "get_text":
                text = ctrl.window_text()
                return f"**'{target}' 的文本内容**:\n{text}"
            elif action == "select":
                if not value:
                    return "select 操作需要提供 value 参数（选项名）"
                ctrl.select(value)
                return f"✅ 已选择 '{target}' 中的 '{value}'"
            elif action == "scroll_down":
                ctrl.scroll("down", "page")
                return f"✅ 已向下滚动 '{target}'"
            elif action == "scroll_up":
                ctrl.scroll("up", "page")
                return f"✅ 已向上滚动 '{target}'"

        except Exception as e:
            return f"UI操作失败: {e}"

    def _tool_app_keyboard(self, args: dict, user_id: str = "") -> str:
        """先聚焦窗口再发送键盘快捷键（pyautogui + pywinauto + 系统媒体键）"""
        keys = args.get("keys", "")
        window_title = args.get("window_title", "")
        wait_before = int(args.get("wait_before", 0))

        if not keys:
            return "请提供要发送的按键，如：'space'、'ctrl+right'"

        try:
            import pyautogui
        except ImportError:
            return "pyautogui 未安装，请执行 pip install pyautogui"

        try:
            import time

            # ---- 等待应用启动 ----
            if wait_before > 0:
                time.sleep(wait_before)

            # ---- 聚焦目标窗口 ----
            if window_title:
                try:
                    from pywinauto import Desktop
                    dlg = None
                    for attempt in range(4):
                        for backend in ["uia", "win32"]:
                            try:
                                desktop = Desktop(backend=backend)
                                for w in desktop.windows():
                                    try:
                                        wt = w.window_text()
                                        if not wt:
                                            continue
                                        wt_lower = wt.lower()
                                        kw = window_title.lower()
                                        if kw in wt_lower or wt_lower in kw:
                                            dlg = w
                                            break
                                        if kw in ("记事本","notepad") and ("记事本" in wt or "notepad" in wt_lower):
                                            dlg = w
                                            break
                                    except Exception:
                                        continue
                                if dlg:
                                    break
                            except Exception:
                                continue
                        if dlg:
                            break
                        time.sleep(0.5)

                    if dlg:
                        dlg.set_focus()
                        time.sleep(0.5)
                    else:
                        # 不阻断——即使没找到窗口也继续发按键
                        pass
                except ImportError:
                    pass

            # ---- 发送按键 ----
            key_str = keys.strip()
            if "+" in key_str:
                combo = key_str.lower().split("+")
                combo = [k.strip() for k in combo if k.strip()]
                pyautogui.hotkey(*combo)
                detail = f"组合键 {'+'.join(combo)}"
            else:
                k = key_str.lower().strip()
                # 系统级媒体键（无需窗口焦点，直接控制任何播放器）
                MEDIA_KEYS = {
                    "media_play_pause": 0xB3,
                    "media_stop": 0xB2,
                    "media_next": 0xB0,
                    "media_prev": 0xB1,
                    "volume_up": 0xAF,
                    "volume_down": 0xAE,
                    "volume_mute": 0xAD,
                }
                if k in MEDIA_KEYS:
                    import ctypes
                    vk = MEDIA_KEYS[k]
                    ctypes.windll.user32.keybd_event(vk, 0, 0, 0)
                    time.sleep(0.05)
                    ctypes.windll.user32.keybd_event(vk, 0, 2, 0)
                    detail = f"系统媒体键 {k}"
                else:
                    pyautogui.press(k)
                    detail = f"按键 {k}"

            target_info = f" → 窗口 '{window_title}'" if window_title else ""
            return f"✅ 已发送{detail}{target_info}"

        except Exception as e:
            return f"键盘操作失败: {e}"

    # ==================== 截图 + 视觉分析 ====================

    def _tool_app_screenshot(self, args: dict, user_id: str = "") -> str:
        """截图并可选调用视觉模型分析（让 AI '看到'屏幕）"""
        try:
            import pyautogui
        except ImportError:
            return "pyautogui 未安装"

        region = args.get("region", "full")  # full | active_window
        question = args.get("question", "")

        win_abs_left = 0  # 窗口在屏幕上的绝对坐标
        win_abs_top = 0
        try:
            if region == "active_window":
                try:
                    from pywinauto import Desktop
                    dlg = Desktop(backend="uia").window(active_only=True)
                    rect = dlg.rectangle()
                    win_abs_left, win_abs_top = rect.left, rect.top
                    img = pyautogui.screenshot(region=(
                        rect.left, rect.top, rect.width(), rect.height()
                    ))
                except Exception:
                    img = pyautogui.screenshot()
            else:
                img = pyautogui.screenshot()

            # 压缩
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=60)
            img_bytes = buf.getvalue()

            if len(img_bytes) > 500 * 1024:
                img = img.resize((img.width // 2, img.height // 2))
                buf = io.BytesIO()
                img.save(buf, format="JPEG", quality=40)
                img_bytes = buf.getvalue()

            b64 = base64.b64encode(img_bytes).decode("utf-8")
            data_url = f"data:image/jpeg;base64,{b64}"

            # 坐标系说明
            if win_abs_left or win_abs_top:
                coord_info = (
                    f"截图区域: {img.width}×{img.height}, {len(img_bytes)//1024}KB。"
                    f"该截图是活动窗口，窗口左上角在屏幕的绝对坐标为 ({win_abs_left}, {win_abs_top})。"
                    f"截图内任意点的屏幕绝对坐标 = (窗口左上角X{win_abs_left} + 截图内X, 窗口左上角Y{win_abs_top} + 截图内Y)。"
                )
            else:
                coord_info = (
                    f"全屏截图: {img.width}×{img.height}, {len(img_bytes)//1024}KB。"
                    f"截图坐标即屏幕绝对坐标，可直接用于 app_click_at。"
                )

            if not question:
                return f"{coord_info}\n图片已生成。如需分析请带 question 参数。"

            if not SILICONFLOW_API_KEY and not ZHIPU_API_KEY:
                return f"{coord_info}\n(视觉分析不可用：请配置 SiliconFlow 或智谱 Key)"

            # 构建坐标定位指令
            coord_instruction = f"截图尺寸 {img.width}×{img.height}。"
            if win_abs_left or win_abs_top:
                coord_instruction += (
                    f"这是活动窗口截图，窗口左上角在屏幕绝对坐标 ({win_abs_left}, {win_abs_top})。"
                    f"返回坐标必须是屏幕绝对坐标！公式：绝对X={win_abs_left}+截图内X，绝对Y={win_abs_top}+截图内Y。"
                )
            else:
                coord_instruction += "这是全屏截图，坐标即屏幕绝对坐标。"

            from openai import OpenAI
            analysis = None

            # 方案1: SiliconFlow Qwen2.5-VL (Key 可用)
            if SILICONFLOW_API_KEY:
                try:
                    client = OpenAI(
                        api_key=SILICONFLOW_API_KEY,
                        base_url=SILICONFLOW_BASE_URL,
                        timeout=30.0,
                    )
                    response = client.chat.completions.create(
                        model="Qwen/Qwen3-VL-8B-Instruct",
                        messages=[{
                            "role": "user",
                            "content": [
                                {"type": "text", "text": f"{coord_instruction}\n\n{question}"},
                                {"type": "image_url", "image_url": {"url": data_url}},
                            ],
                        }],
                        temperature=0.1,
                        max_tokens=500,
                    )
                    analysis = response.choices[0].message.content
                except Exception as e:
                    logger.warning(f"[screenshot] SiliconFlow 视觉分析失败: {e}")

            # 方案2: 智谱 glm-4v-flash 兜底
            if not analysis and ZHIPU_API_KEY:
                try:
                    client = OpenAI(
                        api_key=ZHIPU_API_KEY,
                        base_url=ZHIPU_BASE_URL,
                        timeout=30.0,
                    )
                    response = client.chat.completions.create(
                        model="glm-4v-flash",
                        messages=[{
                            "role": "user",
                            "content": [
                                {"type": "text", "text": f"{coord_instruction}\n\n{question}"},
                                {"type": "image_url", "image_url": {"url": data_url}},
                            ],
                        }],
                        temperature=0.1,
                        max_tokens=500,
                    )
                    analysis = response.choices[0].message.content
                except Exception as e:
                    logger.warning(f"[screenshot] 智谱视觉分析也失败: {e}")

            if analysis:
                return f"{coord_info}\n\n**视觉分析结果**:\n{analysis}"
            return f"{coord_info}\n(视觉分析失败：所有视觉模型不可用)"

        except Exception as e:
            return f"截图失败: {e}"

    # ==================== 坐标点击 ====================

    def _tool_app_click_at(self, args: dict, user_id: str = "") -> str:
        """在屏幕指定坐标点击"""
        try:
            import pyautogui
        except ImportError:
            return "pyautogui 未安装"

        x = args.get("x")
        y = args.get("y")
        button = args.get("button", "left")
        clicks = int(args.get("clicks", 1))
        interval = float(args.get("interval", 0.2))

        if x is None or y is None:
            return "请提供点击坐标 x 和 y。使用 app_screenshot + 视觉分析来确定坐标。"

        try:
            pyautogui.click(int(x), int(y), clicks=clicks, interval=interval, button=button)
            return f"✅ 已{button}键点击坐标 ({x}, {y})，{clicks}次"
        except Exception as e:
            return f"点击失败: {e}"

    # ==================== 剪贴板 ====================

    def _tool_app_clipboard(self, args: dict, user_id: str = "") -> str:
        """读取/写入剪贴板 + 复制粘贴快捷键"""
        action = args.get("action", "paste")
        text = args.get("text", "")

        try:
            if action == "copy":
                import pyautogui
                pyautogui.hotkey("ctrl", "c")
                time.sleep(0.2)
                # 读取剪贴板内容
                try:
                    result = subprocess.run(
                        ["powershell", "-command", "Get-Clipboard"],
                        capture_output=True, text=True, timeout=5,
                    )
                    clip_text = result.stdout.strip()
                    return f"✅ 已复制到剪贴板: {clip_text[:200]}"
                except Exception:
                    return "✅ 已按 Ctrl+C 复制"

            elif action == "paste":
                import pyautogui
                pyautogui.hotkey("ctrl", "v")
                return "✅ 已粘贴 (Ctrl+V)"

            elif action == "get":
                try:
                    result = subprocess.run(
                        ["powershell", "-command", "Get-Clipboard"],
                        capture_output=True, text=True, timeout=5,
                    )
                    clip_text = result.stdout.strip()
                    if not clip_text:
                        return "剪贴板为空"
                    return f"剪贴板内容:\n{clip_text[:2000]}"
                except Exception as e:
                    return f"读取剪贴板失败: {e}"

            elif action == "set":
                if not text:
                    return "请提供要写入剪贴板的 text"
                # 用 clip.exe
                try:
                    proc = subprocess.run(
                        ["clip"], input=text, text=True, timeout=5,
                    )
                    if proc.returncode == 0:
                        return f"✅ 已写入剪贴板 ({len(text)}字)"
                except Exception:
                    pass
                # 备用：PowerShell
                try:
                    escaped = text.replace('"', '`"')
                    subprocess.run(
                        ["powershell", "-command", f'Set-Clipboard -Value "{escaped}"'],
                        timeout=5,
                    )
                    return f"✅ 已写入剪贴板 ({len(text)}字)"
                except Exception as e:
                    return f"写入剪贴板失败: {e}"
            else:
                return f"不支持的操作: {action}。支持: copy, paste, get, set"
        except Exception as e:
            return f"剪贴板操作失败: {e}"

    # ==================== 窗口列表 ====================

    def _tool_app_list_windows(self, args: dict, user_id: str = "") -> str:
        """列出所有打开的窗口"""
        try:
            from pywinauto import Desktop
        except ImportError:
            return "pywinauto 未安装"

        try:
            lines = ["**当前打开的窗口**\n"]
            for backend in ["uia", "win32"]:
                try:
                    desktop = Desktop(backend=backend)
                    for i, w in enumerate(desktop.windows()):
                        try:
                            wt = w.window_text()
                            if not wt or len(wt) < 2:
                                continue
                            if wt.startswith("GDI+") or wt == "Default IME":
                                continue
                            try:
                                rect = w.rectangle()
                                size = f"{rect.width()}×{rect.height()}"
                            except Exception:
                                size = "?"
                            lines.append(f"  {i+1}. {wt[:80]} ({size})")
                        except Exception:
                            continue
                    break  # UIA 成功就不试 win32
                except Exception:
                    continue

            if len(lines) == 1:
                return "未检测到任何窗口"

            return "\n".join(lines[:30])

        except Exception as e:
            return f"获取窗口列表失败: {e}"

    # ==================== 鼠标拖拽 ====================

    def _tool_app_drag(self, args: dict, user_id: str = "") -> str:
        """鼠标拖拽操作"""
        try:
            import pyautogui
        except ImportError:
            return "pyautogui 未安装"

        x1 = args.get("x1")
        y1 = args.get("y1")
        x2 = args.get("x2")
        y2 = args.get("y2")
        duration = float(args.get("duration", 0.5))

        if None in (x1, y1, x2, y2):
            return "请提供起始坐标 (x1, y1) 和目标坐标 (x2, y2)"

        try:
            pyautogui.moveTo(int(x1), int(y1))
            pyautogui.drag(int(x2) - int(x1), int(y2) - int(y1), duration=duration)
            return f"✅ 已从 ({x1},{y1}) 拖拽到 ({x2},{y2})"
        except Exception as e:
            return f"拖拽失败: {e}"

    # ==================== 复合工具：打开+输入 ====================

    def _tool_app_write_text(self, args: dict, user_id: str = "") -> str:
        """
        一站式：打开应用(如未打开)→等待加载→聚焦窗口→输入文字
        解决 LLM 只打开不输入的顽疾
        """
        app_name = args.get("app", args.get("app_name", ""))
        text = args.get("text", args.get("content", ""))
        window_title = args.get("window_title", app_name)  # 窗口标题默认同应用名

        if not app_name:
            return "请提供要打开的应用名称(app)和要输入的文字(text)"
        if not text:
            return "请提供要输入的文字(text)"

        # Step 1: 检查应用是否已打开
        already_open = False
        try:
            from pywinauto import Desktop
            for backend in ["uia", "win32"]:
                try:
                    for w in Desktop(backend=backend).windows():
                        try:
                            wt = w.window_text()
                            if wt and (app_name.lower() in wt.lower() or
                                       window_title.lower() in wt.lower()):
                                already_open = True
                                break
                        except Exception:
                            continue
                except Exception:
                    continue
        except Exception:
            pass

        # Step 2: 打开应用（如果未开）
        if not already_open:
            open_result = self._tool_open_file({"path": app_name}, user_id)
            if "✅" not in open_result:
                return f"无法打开 {app_name}: {open_result}"
            time.sleep(2)  # 等应用加载

        # Step 3: 输入文字
        time.sleep(1)
        return self._do_type_text(text, window_title)

    @staticmethod
    def _do_type_text(text: str, window_title: str = "") -> str:
        """输入文字：剪贴板写入→确保窗口聚焦→点击→粘贴"""
        import pyautogui

        # 1. 写入剪贴板
        clip_ok = False
        try:
            subprocess.run(["clip"], input=text, text=True, timeout=5)
            clip_ok = True
        except Exception:
            try:
                escaped = text.replace('"', '`"')
                subprocess.run(
                    ["powershell", "-command", f'Set-Clipboard -Value "{escaped}"'],
                    timeout=5,
                )
                clip_ok = True
            except Exception:
                pass

        if not clip_ok:
            # 兜底：直接用 pyautogui 打字
            try:
                pyautogui.typewrite(text, interval=0.03)
                return f"✅ 已输入文字 ({len(text)}字)"
            except Exception as e:
                return f"输入失败: {e}"

        # 2. 确保窗口聚焦 + 点击激活编辑区
        if window_title:
            try:
                from pywinauto import Desktop
                for backend in ["uia", "win32"]:
                    try:
                        desktop = Desktop(backend=backend)
                        for w in desktop.windows():
                            try:
                                wt = w.window_text()
                                if wt and window_title.lower() in wt.lower():
                                    w.set_focus()
                                    time.sleep(0.3)
                                    # 点击窗口中间偏上区域（通常是编辑区）
                                    rect = w.rectangle()
                                    cx = rect.left + rect.width() // 2
                                    cy = rect.top + rect.height() // 3
                                    pyautogui.click(cx, cy)
                                    time.sleep(0.2)
                                    break
                            except Exception:
                                continue
                    except Exception:
                        continue
            except Exception:
                pass

        # 3. Ctrl+V 粘贴（发两次确保生效）
        pyautogui.hotkey("ctrl", "v")
        time.sleep(0.15)
        pyautogui.hotkey("ctrl", "v")
        time.sleep(0.1)

        return f"✅ 已粘贴文字到窗口 ({len(text)}字)"

    # ==================== MCP 工具代理 ====================

    async def _execute_mcp_tool(self, tool_name: str, args: dict) -> dict:
        """执行 MCP 外部工具——将调用转发到外部 MCP Server"""
        try:
            from app.mcp_client import get_mcp_manager
            manager = await get_mcp_manager()
            result = await manager.call_tool(tool_name, args)
            logger.info(f"MCP工具结果 [{tool_name}]: {str(result)[:200]}")
            return {"content": str(result), "need_confirm": False}
        except Exception as e:
            logger.error(f"MCP工具异常 [{tool_name}]: {e}")
            return {"content": f"MCP工具执行出错: {e}", "need_confirm": False}

    # ==================== 路径处理 ====================

    def _resolve_path(self, path: str) -> str:
        """解析路径，存在则返回真实路径，不存在返回空串"""
        if not path:
            return ""
        p = os.path.expanduser(path)
        p = os.path.expandvars(p)
        p = os.path.normpath(p)
        if os.path.exists(p):
            return p
        return ""

    def _is_protected(self, path: str) -> bool:
        """检查是否是受保护的系统路径"""
        p = os.path.normpath(path).lower()
        for prot in PROTECTED_PATHS:
            if p.startswith(prot.lower()):
                return True
        return False

    # ==================== 图片分析 ====================

    def _tool_analyze_image(self, args: dict, user_id: str = "") -> str:
        """读取图片文件并使用视觉模型分析描述"""
        path = args.get("path", args.get("file", ""))
        question = args.get("question", "请详细描述这张图片的内容")
        if not path:
            return "请提供图片文件路径"

        real = self._resolve_path(path)
        if not real:
            return f"文件不存在: {path}"

        # 检查文件大小（限制 10MB）
        if os.path.getsize(real) > 10 * 1024 * 1024:
            return f"图片太大（>10MB）: {path}"

        # 读取图片并转 base64
        try:
            ext = Path(real).suffix.lower()
            mime_map = {
                ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
                ".gif": "image/gif", ".bmp": "image/bmp", ".webp": "image/webp",
            }
            mime_type = mime_map.get(ext, "image/jpeg")

            with open(real, "rb") as f:
                img_data = f.read()
            b64 = base64.b64encode(img_data).decode("utf-8")
            data_url = f"data:{mime_type};base64,{b64}"
        except Exception as e:
            return f"读取图片失败: {e}"

        # 调用视觉模型（优先 SiliconFlow Qwen2.5-VL，智谱兜底）
        try:
            from openai import OpenAI

            guard_prompt = (
                "请严格描述这张图片中你实际看到的内容（场景、物体、颜色、文字、人物特征等），"
                "不要猜测或识别具体人物身份。如果图片中有文字，请逐字逐句完整抄录。"
                f"用户问题：{question}"
            )

            # 方案1: SiliconFlow (Key 已验证可用)
            if SILICONFLOW_API_KEY:
                try:
                    client = OpenAI(
                        api_key=SILICONFLOW_API_KEY,
                        base_url=SILICONFLOW_BASE_URL,
                        timeout=30.0,
                    )
                    response = client.chat.completions.create(
                        model="Qwen/Qwen3-VL-8B-Instruct",
                        messages=[{
                            "role": "user",
                            "content": [
                                {"type": "text", "text": guard_prompt},
                                {"type": "image_url", "image_url": {"url": data_url}},
                            ],
                        }],
                        temperature=0.1,
                        max_tokens=500,
                    )
                    return f"**图片分析** ({path})\n\n{response.choices[0].message.content}"
                except Exception as e:
                    logger.warning(f"[analyze_image] SiliconFlow 失败: {e}，尝试智谱")

            # 方案2: 智谱 glm-4v-flash
            if ZHIPU_API_KEY:
                client = OpenAI(
                    api_key=ZHIPU_API_KEY,
                    base_url=ZHIPU_BASE_URL,
                    timeout=30.0,
                )
                response = client.chat.completions.create(
                    model="glm-4v-flash",
                    messages=[{
                        "role": "user",
                        "content": [
                            {"type": "text", "text": guard_prompt},
                            {"type": "image_url", "image_url": {"url": data_url}},
                        ],
                    }],
                    temperature=0.1,
                    max_tokens=500,
                )
                return f"**图片分析** ({path})\n\n{response.choices[0].message.content}"

            return "图片分析不可用（请配置 SiliconFlow 或智谱 API Key）"
        except Exception as e:
            return f"图片分析失败: {e}"

    # ==================== 翻译 ====================

    async def _tool_translate(self, args: dict, user_id: str = "") -> str:
        """使用智谱 GLM 翻译文本（未配置时回退到主模型）"""
        text = args.get("text", "")
        target_lang = args.get("target_lang", args.get("target", "中文"))
        source_lang = args.get("source_lang", args.get("source", ""))

        if not text:
            return "请提供要翻译的文本"

        source_hint = f"从{source_lang}" if source_lang else ""
        prompt = (
            f"你是一位专业翻译。请将以下文本{source_hint}翻译为**{target_lang}**。\n"
            f"要求：\n"
            f"1. 翻译准确、自然、符合{target_lang}表达习惯\n"
            f"2. 保留原文的格式（如列表、加粗等）\n"
            f"3. 只输出翻译结果，不要解释\n\n"
            f"原文：\n{text}"
        )

        try:
            client = self._zhipu_client
            if client is None:
                return "翻译服务不可用（未配置 LLM 客户端）"

            response = await client.ainvoke(
                [{"role": "user", "content": prompt}],
            )
            result = response.content or "翻译失败"
            return f"**翻译结果** ({source_lang or '自动检测'} → {target_lang})\n\n{result}"
        except Exception as e:
            return f"翻译失败: {e}"
